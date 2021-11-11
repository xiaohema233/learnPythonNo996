# -*- coding: utf-8 -*-
"""
Created on Sat Aug 29 19:33:17 2020

@author: Administrator
"""
import xlsxwriter
from pptx import Presentation
myPrs=Presentation(r"H:\示例\第7章\分析报告.pptx")
wb = xlsxwriter.Workbook(r"H:\示例\第7章\报告提取.xlsx")
for i, e in enumerate(myPrs.slides):
    title=e.placeholders[0].text
    profile=e.placeholders[2].text
    honor=e.placeholders[16].text
    ws = wb.add_worksheet(title.replace("财务分析报告",""))
    table = e.placeholders[14].table
    rs,cs=len(table.rows),len(table.columns)
    for r in range(rs):
        for c in range(cs):
            txt=table.cell(r, c).text_frame.paragraphs[0].text
            ws.write(r, c, txt)
    ws.write(r+1, 0, title)	
    ws.write(r+2, 0, profile)	
    ws.write(r+3, 0, honor)	
wb.close()


