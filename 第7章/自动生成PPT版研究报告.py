import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.util import Inches,Pt
myPrs = Presentation('H:\示例\第7章\报告模板.pptx')
xlfile='H:\示例\第7章\公司数据.xlsx'
df0 = pd.read_excel(xlfile,sheet_name=None,header=None)
companys=["上海公司","四川公司","重庆公司","深圳公司"]
for i, e in enumerate(companys):
    df = df0[e]
    data_table=np.array(df.iloc[:9,:]).tolist()
    data_chart=np.array(df.iloc[[0,9],1:]).tolist()
    companyName=df.iat[10,1]
    companyProfile=df.iat[11,1]
    companyHonor=df.iat[12,1]
    slide = myPrs.slides.add_slide(myPrs.slide_layouts[12])
    slide.placeholders[0].text=companyName+'财务分析报告'
    slide.placeholders[2].text=companyProfile
    slide.placeholders[16].text=companyHonor
    slide.placeholders[17].text="公司简介"
    slide.placeholders[18].text="经营范围"
    slide.placeholders[19].text="主要财务指标一览表"
    slide.placeholders[20].text="近五年存货周转天数图示（单位：天）"
    slide.placeholders[22].text="第"+str(i+1)+"页"
    rows,cols=9,6
    pgf = slide.placeholders[14].insert_table(rows,cols)
    tb = pgf.table
    tbl =  pgf.element.graphic.graphicData.tbl
    tbl.tblPr[0].text='{5940675A-B579-460E-94D1-54222C63F5DA}'
    tb.columns[0].width= Inches(1.7)  
    for c in range(1,6):
        tb.columns[c].width= Inches(0.7)  
    for r in range(0,9):
        tb.rows[r].height= Inches(0.32)  
    for row in range(rows):
        for col in range(cols):
            tb.cell(row, col).text_frame.clear()
            new = tb.cell(row, col).text_frame.paragraphs[0]
            new.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            new.text = str(data_table[row][col])
            new.font.size = Pt(10)
            new.font.name =  '微软雅黑'
            new.font.color.rgb =  RGBColor(0, 0, 0)
            new.font.bold = True 
    chart_data = ChartData()
    chart_data.categories = data_chart[0]
    chart_data.add_series('Series 1', data_chart[1])
    pgf=slide.placeholders[15].insert_chart(51, chart_data)
    chart = pgf.chart
    chart.chart_style = 1 
    category_axis = chart.category_axis 
    category_axis.has_major_gridlines = False 
    category_axis.tick_labels.font.italic = True  
    category_axis.tick_labels.font.size = Pt(15)  
    category_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0) 
    value_axis = chart.value_axis 
    value_axis.maximum_scale = 100.0  
    value_axis.minimum_scale = 0.0  
    value_axis.has_major_gridlines = False  
    value_axis.tick_labels.font.size = Pt(10) 
    value_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)  
    plot = chart.plots[0] 
    plot.has_data_labels = True  
    data_labels = plot.data_labels 
    data_labels.font.size = Pt(15)  
    data_labels.font.bold = True
    data_labels.font.color.rgb = RGBColor(0, 0, 0)  
myPrs.save('H:\示例\第7章\分析报告.pptx')