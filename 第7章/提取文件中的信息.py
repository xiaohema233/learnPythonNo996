# -*- coding: utf-8 -*-
"""
Created on Sat Aug 29 19:33:17 2020

@author: Administrator
"""
 
from pptx import Presentation
myPrs = Presentation(r"H:\示例\第7章\提取案例.pptx ")
for slide in myPrs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(shape.text)
        if shape.has_table:
            table = shape.table
            rs,cs=len(table.rows), len(table.columns)
            for r in range(rs):
                data = [table.cell(r, c).text for c in range(cs)]
                print(data)
        if shape.has_chart:
            chart = shape.chart
            for plot in chart.plots:
                category_labels = [c.label for c in plot.categories]
                print(category_labels)
                for serie in plot.series:
                    print(serie.values)

